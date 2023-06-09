from file_parser.parser_interface import Parser
from pptx import Presentation
from io import BytesIO
from file_parser.utils import singleton

@singleton
class PPTXParser(Parser):

    def parser_buffer(self, buffer):
        prs = Presentation(BytesIO(buffer))
        result = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    result.append(paragraph.text)
        return result
    
    def parse(self, buffer, length = 5, cond_split="\n\n\n") -> list:
        content = self.parser_buffer(buffer)
        print(content)
        paragragh_list = [paragragh for paragragh in content if len(paragragh) > length]
        return paragragh_list